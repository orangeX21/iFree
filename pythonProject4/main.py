# import tkinter as tk
# from tkinter import filedialog
# import pptx
#
# window = tk.Tk()
# window.title('PPT批量处理')
# window.geometry('500x300')
# # 设置窗口最小尺寸
# window.minsize(400, 300)
#
# # 设置窗口最大尺寸
# window.maxsize(800, 600)
#
# ppt_file = ''
# save_file = ''
#
#
# def select_file():
#     global ppt_file, save_file
#     ppt_file = filedialog.askopenfilename(filetypes=[('PPT Files', '*.pptx')])
#     save_file = filedialog.asksaveasfilename(defaultextension='.pptx')
#     filename_label['text'] = ppt_file
#     process_ppt(ppt_file, save_file)
#
#
# def process_ppt(ppt_file, save_file):
#     prs = pptx.Presentation(ppt_file)
#     for slide in prs.slides:
#         for shape in slide.shapes:
#             if shape.has_text_frame:
#                 text_frame = shape.text_frame
#
#                 # 调整文字大小
#                 text_frame.text = text_frame.text.upper()
#     prs.save(save_file)
#     label = tk.Label(window, text='处理完成!')
#     label.pack()
#
#
# filename_label = tk.Label(window, text='未选择文件')
# filename_label.pack()
#
# btn1 = tk.Button(text='选择文件', command=select_files)
# btn2 = tk.Button(text='开始处理', command=process_files)
#
# btn1.pack()
# btn2.pack()
# window.mainloop()

import tkinter as tk
from tkinter import filedialog, messagebox
import pptx
import os

# import click_exit
# import frame
# import process

window = tk.Tk()
window.title('调整PPT文字格式的脚本')
window.geometry('500x300')

window.minsize(400, 300)

# 设置窗口最大尺寸
window.maxsize(800, 600)

ppt_file = ''
save_file = ''
filename = tk.StringVar()
filename.set('未选择文件')
save_dir = ''  # 定义为全局变量


def select_files():
    global ppt_file, save_file, save_dir
    # ppt_file = filedialog.askopenfilename()
    ppt_file = filedialog.askopenfilename()

    # 检查是否选择了文件
    if ppt_file:
        dir_name = os.path.dirname(ppt_file)
        save_dir = os.path.join(dir_name, 'modified')
        # os.mkdir(save_dir)
        if not os.path.exists(save_dir):
            os.mkdir(save_dir)

    else:

        # 未选择,打印提示
        messagebox.showinfo('开始', '请先选择文件!')

    # save_file = filedialog.asksaveasfilename()
    save_file = os.path.join(save_dir, 'new_' + os.path.basename(ppt_file))
    files = ppt_file
    file_list = []
    # for file in files:
    #     file_list.append(file)

    # 拼接文件列表并显示
    if isinstance(ppt_file, list):
        file_list = ppt_file
    else:
        file_list = [ppt_file]
    filename.set('\n'.join(file_list))


def click_exit():
    window.quit()


def process_files():
    if not ppt_file:
        messagebox.showwarning('警告', '请先选择文件!')
    else:
        if ppt_file and save_file:
            prs = pptx.Presentation(ppt_file)
            #
            # # 获取首文本框格式作为样板
            # src_shape = prs.slides[0].shapes[0]
            # src_font = src_shape.text_frame.paragraphs[0].font
            #
            # for slide in prs.slides:
            #
            #     for shape in slide.shapes:
            #
            #         if shape.has_text_frame:
            #             text_frame = shape.text_frame
            #
            #             # 设置字号和高度
            #             text_frame.paragraphs[0].font.size = src_font.size
            #             shape.height = src_shape.height
            #
            # # 保存为save_file
            # prs.save(save_file)

            # # 处理PPT逻辑
            #
            prs.save(save_file)

            messagebox.showinfo('完成', f'已成功保存到\n{save_file}')
        else:
            messagebox.showwarning('警告', '请先选择文件!')


# process.process_files()
messagebox.showinfo('开始', '请先选择文件!')
# frame.frame()
file_frame = tk.Frame(
    window,
    bg='yellow',
    relief='groove',

    height=100,  # 高度设大
    width=100,  # 宽度设大
    padx=20,
    pady=20
)
file_label = tk.Label(file_frame, font=('Arial', 12), textvariable=filename)
file_label.pack(fill='x')
file_frame.pack(fill='x')
btn1 = tk.Button(window, text='选择文件', command=select_files)
btn1.pack(anchor='n')
btn2 = tk.Button(window, text='开始处理', command=process_files)
btn2.pack(anchor='n')
exit_btn = tk.Button(text="退出", command=click_exit)
exit_btn.pack(anchor='n')
window.mainloop()

#
# select_btn = tk.Button(text="选择文件")
# select_btn.pack(side="left")
#
# # 退出按钮
# exit_btn = tk.Button(text="退出")
# exit_btn.pack() # 不设置side
#
# # 开始处理按钮
# process_btn = tk.Button(text="开始处理")
# process_btn.pack(side="right")
#
# # 其他窗口组件、布局代码
#
# # 为按钮添加点击事件
# exit_btn.config(command=window.quit)
